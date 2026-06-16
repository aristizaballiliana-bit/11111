"""
ingest.py - 轻量版
用 SQLite + numpy 做向量存储和检索，替换 ChromaDB。
内存占用约 50-80MB，适合 Railway 免费套餐。
"""

import os
import io
import json
import sqlite3
import numpy as np
import voyageai
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

VOYAGE_API_KEY = os.environ.get("VOYAGE_API_KEY")
DB_PATH = os.environ.get("CHROMA_PATH", "./chroma_db").rstrip("/") + "/kb.sqlite3"

vo = None
if VOYAGE_API_KEY:
    vo = voyageai.Client(api_key=VOYAGE_API_KEY)


def _get_conn():
    os.makedirs(os.path.dirname(DB_PATH), exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source TEXT NOT NULL,
            text TEXT NOT NULL,
            embedding TEXT NOT NULL
        )
    """)
    conn.commit()
    return conn


# ---------- 文件解析 ----------

def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    for table in doc.tables:
        for row in table.rows:
            text += "\n" + " | ".join(cell.text for cell in row.cells)
    return text


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"【工作表: {sheet.title}】")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_texts.append(line.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                slide_texts.append(f"[备注] {notes.strip()}")
        if slide_texts:
            parts.append(f"【幻灯片 {i}】\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_txt(file_bytes: bytes) -> str:
    for enc in ("utf-8", "gbk", "gb2312"):
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text_any(filename: str, file_bytes: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif name.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    elif name.endswith((".xlsx", ".xlsm")):
        return extract_text_from_xlsx(file_bytes)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    elif name.endswith(".txt"):
        return extract_text_from_txt(file_bytes)
    else:
        raise ValueError("仅支持 .pdf、.docx、.xlsx、.pptx、.txt 文件")


# ---------- 向量化 ----------

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list:
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        chunk = text[start:start + chunk_size]
        if chunk.strip():
            chunks.append(chunk.strip())
        start += chunk_size - overlap
    return chunks


def embed_texts(texts: list) -> list:
    if not vo:
        raise RuntimeError("VOYAGE_API_KEY 未配置")
    result = vo.embed(texts, model="voyage-3", input_type="document")
    return result.embeddings


# ---------- 知识库操作 ----------

def ingest_document(filename: str, file_bytes: bytes, metadata: dict = None) -> dict:
    text = extract_text_any(filename, file_bytes)
    chunks = chunk_text(text)
    if not chunks:
        return {"filename": filename, "chunks_added": 0, "message": "未提取到文本内容"}

    embeddings = embed_texts(chunks)

    conn = _get_conn()
    conn.execute("DELETE FROM chunks WHERE source = ?", (filename,))
    conn.executemany(
        "INSERT INTO chunks (source, text, embedding) VALUES (?, ?, ?)",
        [(filename, chunk, json.dumps(emb)) for chunk, emb in zip(chunks, embeddings)]
    )
    conn.commit()
    conn.close()

    return {"filename": filename, "chunks_added": len(chunks)}


def query_knowledge_base(query: str, top_k: int = 5) -> list:
    if not vo:
        raise RuntimeError("VOYAGE_API_KEY 未配置")

    query_emb = np.array(
        vo.embed([query], model="voyage-3", input_type="query").embeddings[0]
    )

    conn = _get_conn()
    rows = conn.execute("SELECT source, text, embedding FROM chunks").fetchall()
    conn.close()

    if not rows:
        return []

    sources = [r[0] for r in rows]
    texts = [r[1] for r in rows]
    matrix = np.array([json.loads(r[2]) for r in rows])

    # 余弦相似度
    norms = np.linalg.norm(matrix, axis=1, keepdims=True)
    norms[norms == 0] = 1
    matrix_norm = matrix / norms
    query_norm = query_emb / (np.linalg.norm(query_emb) or 1)
    scores = matrix_norm @ query_norm

    top_indices = np.argsort(scores)[::-1][:top_k]
    return [{"text": texts[i], "source": sources[i]} for i in top_indices]


def list_documents() -> list:
    conn = _get_conn()
    rows = conn.execute("SELECT DISTINCT source FROM chunks").fetchall()
    conn.close()
    return [r[0] for r in rows]


def delete_document(filename: str) -> dict:
    conn = _get_conn()
    cursor = conn.execute("DELETE FROM chunks WHERE source = ?", (filename,))
    deleted = cursor.rowcount
    conn.commit()
    conn.close()
    return {"filename": filename, "deleted_chunks": deleted}
